"""
Standalone PPT generation test without AWS or Strands dependencies.
"""
import os
import sys
from datetime import datetime
from typing import List, Dict, Any

# Add project paths
project_root = os.path.dirname(__file__)
sys.path.insert(0, project_root)
sys.path.insert(0, os.path.join(project_root, 'src'))

def create_standalone_ppt_generator():
    """Create a standalone PPT generator without dependencies."""
    
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        
        class StandalonePPTGenerator:
            def __init__(self):
                self.current_colors = {
                    'primary': RGBColor(30, 58, 138),     # Navy blue
                    'secondary': RGBColor(245, 158, 11),   # Gold
                    'text': RGBColor(31, 41, 55),          # Dark gray
                }
                self.current_fonts = {
                    'title': {'size': Pt(36)},
                    'subtitle': {'size': Pt(24)},
                    'content': {'size': Pt(18)}
                }
            
            def generate_business_presentation(self, company_name: str, use_cases: List[Dict] = None) -> str:
                """Generate a complete business transformation presentation."""
                
                if not use_cases:
                    use_cases = self._get_default_use_cases()
                
                # Create presentation
                ppt = Presentation()
                
                # Title slide
                self._create_title_slide(ppt, company_name)
                
                # Overview slide
                self._create_overview_slide(ppt, company_name, len(use_cases))
                
                # Use case slides (minimum 3)
                for i, use_case in enumerate(use_cases[:3], 1):
                    self._create_use_case_slide(ppt, i, use_case)
                
                # Implementation slide
                self._create_implementation_slide(ppt)
                
                # Next steps slide
                self._create_next_steps_slide(ppt)
                
                # Save presentation
                timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
                filename = f"{company_name.replace(' ', '_').lower()}_transformation_{timestamp}.pptx"
                
                ppt.save(filename)
                return filename
            
            def _create_title_slide(self, ppt: Presentation, company_name: str):
                """Create title slide."""
                title_layout = ppt.slide_layouts[0]
                slide = ppt.slides.add_slide(title_layout)
                
                title = slide.shapes.title
                subtitle = slide.placeholders[1]
                
                title.text = f"Digital Transformation Strategy for {company_name}"
                subtitle.text = "AI-Powered Business Solutions for Competitive Advantage"
                
                # Style title
                title.text_frame.paragraphs[0].font.size = self.current_fonts['title']['size']
                title.text_frame.paragraphs[0].font.color.rgb = self.current_colors['primary']
                title.text_frame.paragraphs[0].font.bold = True
                
                # Style subtitle
                subtitle.text_frame.paragraphs[0].font.size = self.current_fonts['subtitle']['size']
                subtitle.text_frame.paragraphs[0].font.color.rgb = self.current_colors['text']
            
            def _create_overview_slide(self, ppt: Presentation, company_name: str, use_case_count: int):
                """Create overview slide."""
                content_layout = ppt.slide_layouts[1]
                slide = ppt.slides.add_slide(content_layout)
                
                title = slide.shapes.title
                content = slide.placeholders[1]
                
                title.text = "Executive Overview"
                title.text_frame.paragraphs[0].font.size = self.current_fonts['title']['size']
                title.text_frame.paragraphs[0].font.color.rgb = self.current_colors['primary']
                title.text_frame.paragraphs[0].font.bold = True
                
                content.text = f"""• {company_name} is positioned for strategic digital transformation
- Identified {use_case_count} high-impact technology initiatives
- Focus on AI-powered automation and operational efficiency
- Conservative ROI projections with realistic implementation timelines
- Phased approach to minimize business disruption"""
                
                # Style content
                for paragraph in content.text_frame.paragraphs:
                    paragraph.font.size = self.current_fonts['content']['size']
                    paragraph.font.color.rgb = self.current_colors['text']
            
            def _create_use_case_slide(self, ppt: Presentation, number: int, use_case: Dict):
                """Create individual use case slide."""
                content_layout = ppt.slide_layouts[1]
                slide = ppt.slides.add_slide(content_layout)
                
                title = slide.shapes.title
                content = slide.placeholders[1]
                
                title.text = f"Use Case {number}: {use_case['title']}"
                title.text_frame.paragraphs[0].font.size = Pt(32)
                title.text_frame.paragraphs[0].font.color.rgb = self.current_colors['primary']
                title.text_frame.paragraphs[0].font.bold = True
                
                content.text = f"""Challenge: {use_case['challenge']}
Solution: {use_case['solution']}
Benefits: {use_case['benefits']}
Timeline: {use_case['timeline']}
Success Metrics: {use_case['metrics']}"""
                
                # Style content
                for paragraph in content.text_frame.paragraphs:
                    paragraph.font.size = self.current_fonts['content']['size']
                    paragraph.font.color.rgb = self.current_colors['text']
            
            def _create_implementation_slide(self, ppt: Presentation):
                """Create implementation roadmap slide."""
                content_layout = ppt.slide_layouts[1]
                slide = ppt.slides.add_slide(content_layout)
                
                title = slide.shapes.title
                content = slide.placeholders[1]
                
                title.text = "Implementation Roadmap"
                title.text_frame.paragraphs[0].font.size = self.current_fonts['title']['size']
                title.text_frame.paragraphs[0].font.color.rgb = self.current_colors['primary']
                title.text_frame.paragraphs[0].font.bold = True
                
                content.text = """• Phase 1: Assessment and Planning (Months 1-2)
- Phase 2: Pilot Implementation (Months 3-6)
- Phase 3: Full Deployment (Months 7-10)
- Phase 4: Optimization and Scaling (Months 11-12)
- Ongoing: Performance monitoring and continuous improvement"""
                
                # Style content
                for paragraph in content.text_frame.paragraphs:
                    paragraph.font.size = self.current_fonts['content']['size']
                    paragraph.font.color.rgb = self.current_colors['text']
            
            def _create_next_steps_slide(self, ppt: Presentation):
                """Create next steps slide."""
                content_layout = ppt.slide_layouts[1]
                slide = ppt.slides.add_slide(content_layout)
                
                title = slide.shapes.title
                content = slide.placeholders[1]
                
                title.text = "Next Steps & Recommendations"
                title.text_frame.paragraphs[0].font.size = self.current_fonts['title']['size']
                title.text_frame.paragraphs[0].font.color.rgb = self.current_colors['primary']
                title.text_frame.paragraphs[0].font.bold = True
                
                content.text = """• Secure executive sponsorship and transformation budget
- Establish project governance and cross-functional team
- Conduct detailed technical and business readiness assessment
- Select technology partners and implementation vendors
- Develop comprehensive change management strategy"""
                
                # Style content
                for paragraph in content.text_frame.paragraphs:
                    paragraph.font.size = self.current_fonts['content']['size']
                    paragraph.font.color.rgb = self.current_colors['text']
            
            def _get_default_use_cases(self) -> List[Dict]:
                """Get default use cases for demonstration."""
                return [
                    {
                        'title': 'Intelligent Process Automation',
                        'challenge': 'Manual processes causing inefficiencies and errors',
                        'solution': 'AI-powered workflow automation for key business processes',
                        'benefits': '40-50% reduction in processing time and 60% fewer errors',
                        'timeline': '4-6 months implementation',
                        'metrics': 'Process efficiency, error rates, employee productivity'
                    },
                    {
                        'title': 'Predictive Analytics Platform',
                        'challenge': 'Reactive decision-making with limited business insights',
                        'solution': 'Advanced analytics platform for predictive intelligence',
                        'benefits': '30% improvement in forecasting accuracy',
                        'timeline': '6-8 months implementation', 
                        'metrics': 'Forecast accuracy, decision speed, business outcomes'
                    },
                    {
                        'title': 'Customer Experience Enhancement',
                        'challenge': 'Limited personalization and customer engagement',
                        'solution': 'AI-driven customer experience optimization platform',
                        'benefits': '25% increase in customer satisfaction scores',
                        'timeline': '5-7 months implementation',
                        'metrics': 'NPS scores, retention rates, engagement metrics'
                    }
                ]
        
        return StandalonePPTGenerator
        
    except ImportError as e:
        print(f"Import error: {e}")
        print("Please install: pip install python-pptx==0.6.21")
        return None

def test_standalone_ppt():
    """Test standalone PPT generation."""
    
    print("🚀 Testing Standalone PPT Generation")
    print("=" * 50)
    
    try:
        # Create generator
        PPTGenerator = create_standalone_ppt_generator()
        if not PPTGenerator:
            return
        
        generator = PPTGenerator()
        
        # Test with multiple companies and styles
        test_companies = [
            "TechCorp Solutions",
            "InnovateMax Industries", 
            "Digital Dynamics Corp"
        ]
        
        results = []
        
        for company in test_companies:
            print(f"\n📋 Generating presentation for {company}...")
            
            filename = generator.generate_business_presentation(company)
            
            if os.path.exists(filename):
                file_size = os.path.getsize(filename)
                results.append({
                    'company': company,
                    'filename': filename,
                    'size': file_size,
                    'path': os.path.abspath(filename)
                })
                
                print(f"✅ Created: {filename}")
                print(f"📁 Location: {os.path.abspath(filename)}")
                print(f"📊 Size: {file_size:,} bytes")
            else:
                print(f"❌ Failed to create presentation for {company}")
        
        # Summary
        print(f"\n🎉 SUMMARY")
        print("=" * 30)
        print(f"✅ Successfully created {len(results)} presentations")
        
        for result in results:
            print(f"\n📎 {result['company']}:")
            print(f"   File: {result['filename']}")
            print(f"   Size: {result['size']:,} bytes")
            print(f"   Path: {result['path']}")
        
        print(f"\n🔓 All files are fully editable in PowerPoint!")
        print(f"🎨 Template: Executive (Navy blue theme)")
        print(f"📋 Content: Title + Overview + 3 Use Cases + Implementation + Next Steps")
        
        return results
        
    except Exception as e:
        print(f"❌ Error: {e}")
        import traceback
        traceback.print_exc()
        return None

if __name__ == "__main__":
    test_standalone_ppt()