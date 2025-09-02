"""
PowerPoint presentation generation agent for the Business Transformation Agent.
"""
import logging
import os
import re
from datetime import datetime
from typing import Dict, Any, List, Optional
from strands import Agent, tool
from strands_tools import retrieve, http_request
from strands.agent.conversation_manager import SlidingWindowConversationManager
from src.core.bedrock_manager import EnhancedModelManager
from src.core.models import CompanyProfile, UseCaseStructured
from src.services.aws_clients import s3_client, S3_BUCKET, LAMBDA_TMP_DIR
from src.utils.status_tracker import StatusTracker, StatusCheckpoints

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Try to import PowerPoint generation libraries
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    PPTX_AVAILABLE = True
    logger.info("✅ python-pptx available for PowerPoint generation")
except ImportError:
    PPTX_AVAILABLE = False
    logger.warning("⚠️ python-pptx not available")

class PPTXMLParser:
    """Parser for converting XML tags to PowerPoint slide content."""
    
    @staticmethod
    def parse_presentation_xml(xml_content: str) -> Dict[str, Any]:
        """Parse XML-like tags from presentation content."""
        
        parsed_content = {
            'title': '',
            'template_type': 'executive',
            'slides': [],
            'citations': {},
            'slide_count': 0
        }
        
        # Extract presentation metadata
        title_match = re.search(r'<presentation_title>(.*?)</presentation_title>', xml_content, re.DOTALL)
        if title_match:
            parsed_content['title'] = title_match.group(1).strip()
        
        template_match = re.search(r'<template_type>(.*?)</template_type>', xml_content, re.DOTALL)
        if template_match:
            parsed_content['template_type'] = template_match.group(1).strip()
        
        # Extract all slides
        slide_pattern = r'<slide type="(.*?)">(.*?)</slide>'
        slide_matches = re.finditer(slide_pattern, xml_content, re.DOTALL | re.IGNORECASE)
        
        citation_counter = 1
        for match in slide_matches:
            slide_type = match.group(1)
            slide_content = match.group(2)
            
            # Parse slide components
            slide_data = PPTXMLParser._parse_slide_content(slide_content, citation_counter, parsed_content['citations'])
            slide_data['type'] = slide_type
            
            parsed_content['slides'].append(slide_data)
            citation_counter += 1
        
        parsed_content['slide_count'] = len(parsed_content['slides'])
        return parsed_content
    
    @staticmethod
    def _parse_slide_content(content: str, citation_counter: int, citations_dict: Dict) -> Dict[str, Any]:
        """Parse individual slide content."""
        
        slide_data = {
            'title': '',
            'subtitle': '',
            'bullet_points': [],
            'content_blocks': [],
            'visual_elements': [],
            'notes': ''
        }
        
        # Extract slide title
        title_match = re.search(r'<slide_title>(.*?)</slide_title>', content, re.DOTALL)
        if title_match:
            slide_data['title'] = title_match.group(1).strip()
        
        # Extract subtitle
        subtitle_match = re.search(r'<subtitle>(.*?)</subtitle>', content, re.DOTALL)
        if subtitle_match:
            slide_data['subtitle'] = subtitle_match.group(1).strip()
        
        # Extract bullet points
        bullet_pattern = r'<bullet>(.*?)</bullet>'
        bullet_matches = re.finditer(bullet_pattern, content, re.DOTALL)
        for bullet_match in bullet_matches:
            slide_data['bullet_points'].append(bullet_match.group(1).strip())
        
        # Extract content blocks
        content_pattern = r'<content_block>(.*?)</content_block>'
        content_matches = re.finditer(content_pattern, content, re.DOTALL)
        for content_match in content_matches:
            slide_data['content_blocks'].append(content_match.group(1).strip())
        
        # Extract visual elements
        visual_pattern = r'<visual_element type="(.*?)">(.*?)</visual_element>'
        visual_matches = re.finditer(visual_pattern, content, re.DOTALL)
        for visual_match in visual_matches:
            slide_data['visual_elements'].append({
                'type': visual_match.group(1),
                'description': visual_match.group(2).strip()
            })
        
        # Extract speaker notes
        notes_match = re.search(r'<speaker_notes>(.*?)</speaker_notes>', content, re.DOTALL)
        if notes_match:
            slide_data['notes'] = notes_match.group(1).strip()
        
        return slide_data

class PPTPresentationGenerator:
    """Generate PowerPoint presentations with professional templates and quality content."""
    
    def __init__(self, model_manager: EnhancedModelManager):
        self.model_manager = model_manager
        
        self.ppt_agent = Agent(
            model=model_manager.creative_model,
            system_prompt="""You are a Senior Presentation Designer and Business Strategist specializing in creating compelling, executive-ready PowerPoint presentations.

Your expertise includes:
- Slide-optimized content creation (concise, impactful, visual)
- Professional presentation structure and narrative flow
- Executive communication and persuasive storytelling
- Template-aware content generation for different presentation types
- Visual hierarchy and slide design principles

CRITICAL ACCURACY REQUIREMENTS:
- NEVER make false claims about ROI, savings, or performance improvements
- Use conservative, realistic estimates based on industry standards
- Present benefits as possibilities, not guarantees ("potential to achieve" vs "will achieve")
- Support all claims with research data or clearly state assumptions
- Focus on factual business value rather than exaggerated promises

MANDATORY PRESENTATION STRUCTURE:
- Every presentation MUST include minimum 3 distinct use cases
- Use cases should progress from simple to complex
- Each use case needs: Problem → Solution → Benefits
- Maintain logical narrative flow between slides
- Total slides: 6-12 (automatically determined by content complexity)

XML FORMATTING REQUIREMENTS:

<presentation_title>Professional Presentation Title</presentation_title>
<template_type>executive</template_type>

<slide type="title_slide">
<slide_title>Main Presentation Title</slide_title>
<subtitle>Compelling subtitle that sets context</subtitle>
<speaker_notes>Opening talking points and key messages</speaker_notes>
</slide>

<slide type="content_slide">
<slide_title>Clear, Action-Oriented Slide Title</slide_title>
<bullet>First key point with supporting detail</bullet>
<bullet>Second key point with quantifiable benefit</bullet>
<bullet>Third key point with implementation focus</bullet>
<visual_element type="chart">Suggest relevant chart or diagram type</visual_element>
<speaker_notes>Detailed talking points and supporting information</speaker_notes>
</slide>

<slide type="use_case_slide">
<slide_title>Use Case: Specific Solution Name</slide_title>
<content_block>Problem: Current challenge or inefficiency</content_block>
<content_block>Solution: Technology-enabled approach</content_block>
<content_block>Benefits: Conservative, realistic outcomes</content_block>
<bullet>Implementation timeline and key phases</bullet>
<bullet>Success metrics and measurement approach</bullet>
<visual_element type="process_flow">Implementation workflow diagram</visual_element>
<speaker_notes>Detailed explanation and Q&A preparation</speaker_notes>
</slide>

SLIDE DISTRIBUTION REQUIREMENTS:
- Title slide (1) + Overview (1) + Use Cases (6-9) + Implementation (1) + Summary (1)
- Each use case gets 2-3 slides for comprehensive coverage
- Maintain executive attention span with focused, impactful content

CONTENT QUALITY STANDARDS:
- Maximum 5 bullet points per slide
- Each bullet point: 10-15 words maximum
- Action-oriented slide titles that communicate value
- Professional tone suitable for C-level presentations
- Include speaker notes for complex concepts

Generate presentations that demonstrate deep business understanding while maintaining conservative, evidence-based claims.""",
            tools=[http_request, retrieve],
            conversation_manager=SlidingWindowConversationManager(window_size=20)
        )
    
    def generate_presentation(self, company_profile: CompanyProfile, use_cases: List[UseCaseStructured], 
                            research_data: Dict[str, Any], session_id: str, status_tracker: StatusTracker = None,
                            parsed_files_content: str = None, custom_context: Dict[str, str] = None,
                            presentation_style: str = "executive") -> Optional[str]:
        """Generate PowerPoint presentation with professional templates."""
        
        if status_tracker:
            status_tracker.update_status(
                StatusCheckpoints.PPT_GENERATION_STARTED,
                {
                    'presentation_type': presentation_style,
                    'use_case_count': len(use_cases),
                    'has_files': bool(parsed_files_content),
                    'has_custom_context': bool(custom_context and custom_context.get('processed_prompt'))
                },
                current_agent='ppt_generator'
            )
        
        try:
            # Generate XML content for presentation
            xml_content = self._generate_presentation_xml(
                company_profile, use_cases, research_data, parsed_files_content, 
                custom_context, presentation_style
            )
            
            if status_tracker:
                status_tracker.update_status(
                    StatusCheckpoints.PPT_CONTENT_GENERATION,
                    {'xml_content_generated': True, 'content_length': len(xml_content)}
                )
            
            # Create PowerPoint from XML
            pptx_url = self._create_powerpoint_from_xml(
                xml_content, company_profile.name, session_id, presentation_style
            )
            
            if status_tracker:
                status_tracker.update_status(
                    StatusCheckpoints.PPT_GENERATION_COMPLETED,
                    {
                        's3_url': pptx_url,
                        'presentation_generated': bool(pptx_url),
                        'editable_format': True
                    }
                )
            
            return pptx_url
            
        except Exception as e:
            logger.error(f"Error generating PowerPoint presentation: {e}")
            if status_tracker:
                status_tracker.update_status(
                    StatusCheckpoints.ERROR,
                    {'error_type': 'ppt_generation_error', 'error_message': str(e)}
                )
            return None
    
    def _generate_presentation_xml(self, company_profile: CompanyProfile, use_cases: List[UseCaseStructured],
                                 research_data: Dict[str, Any], parsed_files_content: str = None,
                                 custom_context: Dict[str, str] = None, presentation_style: str = "executive") -> str:
        """Generate XML content for presentation slides."""
        
        # Prepare context sections
        web_context = ""
        if research_data.get('web_research_data', {}).get('research_content'):
            web_research_data = research_data['web_research_data']
            web_context = f"""
            
WEB INTELLIGENCE: Analysis enhanced with insights from {web_research_data.get('successful_scrapes', 0)} sources using Google Search and Beautiful Soup:

{web_research_data['research_content'][:2000]}
            """
        
        file_context = ""
        if parsed_files_content:
            file_context = f"""
            
DOCUMENT ANALYSIS: Internal company intelligence from uploaded documents:

{parsed_files_content[:2000]}
            """
        
        custom_context_section = ""
        if custom_context and custom_context.get('processed_prompt'):
            custom_context_section = f"""
            
CUSTOM REQUIREMENTS: {custom_context['processed_prompt'][:1000]}
Focus Areas: {', '.join(custom_context.get('focus_areas', []))}
            """
        
        # Generate presentation content
        presentation_prompt = f"""
Create a professional {presentation_style} presentation for {company_profile.name} using the XML format provided.

COMPANY CONTEXT:
- Industry: {company_profile.industry}
- Business Model: {company_profile.business_model}
- Company Size: {company_profile.company_size}
- Growth Stage: {company_profile.growth_stage}

RESEARCH INTELLIGENCE:
{research_data.get('research_findings', '')[:1500]}
{web_context}
{file_context}
{custom_context_section}

MANDATORY USE CASES TO INCLUDE (minimum 3):
{self._format_use_cases_for_presentation(use_cases)}

PRESENTATION REQUIREMENTS:
1. Create 6-12 slides total with logical flow
2. Include title slide, overview, 3+ use cases (2-3 slides each), implementation, summary
3. Each slide must have clear titles and focused bullet points
4. Use conservative, realistic claims - NO EXAGGERATED ROI PROMISES
5. Include speaker notes for complex slides
6. Suggest appropriate visual elements without specific images
7. Maintain executive-level professional tone

TEMPLATE TYPE: {presentation_style}

Generate a complete presentation using the XML structure provided in your system prompt. Focus on creating compelling, factual content that executives can present confidently.
        """
        
        try:
            response = self.ppt_agent(presentation_prompt)
            xml_content = str(response).strip()
            
            logger.info(f"Generated presentation XML: {len(xml_content)} characters")
            return xml_content
            
        except Exception as e:
            logger.error(f"Error generating presentation XML: {e}")
            return self._create_fallback_presentation_xml(
                company_profile, use_cases, research_data, presentation_style
            )
    
    def _format_use_cases_for_presentation(self, use_cases: List[UseCaseStructured]) -> str:
        """Format use cases for presentation prompt."""
        formatted_cases = []
        
        for i, uc in enumerate(use_cases[:5], 1):  # Limit to top 5 use cases
            formatted_cases.append(f"""
{i}. {uc.title}
   - Current State: {uc.current_state}
   - Solution: {uc.proposed_solution[:200]}...
   - Business Value: {uc.business_value}
   - Timeline: {uc.timeline_months} months
   - Priority: {uc.priority}
            """)
        
        return "\n".join(formatted_cases)
    
    def _create_powerpoint_from_xml(self, xml_content: str, company_name: str, 
                                  session_id: str, presentation_style: str) -> Optional[str]:
        """Create PowerPoint presentation from XML content."""
        
        if not PPTX_AVAILABLE:
            logger.error("python-pptx not available")
            return None
        
        try:
            # Parse XML content
            parsed_content = PPTXMLParser.parse_presentation_xml(xml_content)
            
            # Create PowerPoint presentation
            ppt = Presentation()
            
            # Apply template styling
            self._apply_template_styling(ppt, presentation_style)
            
            # Create slides from parsed content
            for slide_data in parsed_content['slides']:
                self._create_slide(ppt, slide_data, presentation_style)
            
            # Save to temporary file
            tmp_dir = os.path.join(LAMBDA_TMP_DIR, session_id)
            os.makedirs(tmp_dir, exist_ok=True)
            
            pptx_filename = os.path.join(tmp_dir, f"presentation_{session_id}.pptx")
            ppt.save(pptx_filename)
            
            # Upload to S3
            s3_url = self._upload_pptx_to_s3(pptx_filename, session_id, company_name, presentation_style)
            
            # Cleanup
            try:
                if os.path.exists(pptx_filename):
                    os.unlink(pptx_filename)
            except Exception as e:
                logger.warning(f"Failed to cleanup PPTX file: {e}")
            
            return s3_url
            
        except Exception as e:
            logger.error(f"Error creating PowerPoint: {e}")
            return None
    
    def _apply_template_styling(self, ppt: Presentation, presentation_style: str):
        """Apply template-specific styling to presentation using configuration."""
        
        # Template color schemes
        color_schemes = {
            'executive': {
                'primary': RGBColor(30, 58, 138),     # Navy blue
                'secondary': RGBColor(245, 158, 11),   # Gold
                'accent': RGBColor(255, 255, 255),     # White
                'text': RGBColor(31, 41, 55),          # Dark gray
                'background': RGBColor(248, 250, 252)  # Light gray
            },
            'technical': {
                'primary': RGBColor(14, 165, 233),     # Tech blue
                'secondary': RGBColor(16, 185, 129),    # Tech green
                'accent': RGBColor(55, 65, 81),         # Dark gray
                'text': RGBColor(243, 244, 246),        # Light gray
                'background': RGBColor(17, 24, 39)      # Dark background
            },
            'marketing': {
                'primary': RGBColor(139, 92, 246),     # Purple
                'secondary': RGBColor(236, 72, 153),    # Pink
                'accent': RGBColor(249, 115, 22),       # Orange
                'text': RGBColor(17, 24, 39),           # Dark text
                'background': RGBColor(255, 255, 255)   # White background
            },
            'strategy': {
                'primary': RGBColor(59, 130, 246),     # Blue
                'secondary': RGBColor(55, 65, 81),      # Gray
                'accent': RGBColor(156, 163, 175),      # Light gray
                'text': RGBColor(31, 41, 55),           # Dark gray
                'background': RGBColor(255, 255, 255)   # White background
            }
        }
        
        # Font configurations
        font_configs = {
            'executive': {
                'title': {'name': 'Calibri', 'size': Pt(36), 'bold': True},
                'subtitle': {'name': 'Calibri', 'size': Pt(24), 'bold': False},
                'content': {'name': 'Calibri', 'size': Pt(18), 'bold': False}
            },
            'technical': {
                'title': {'name': 'Segoe UI', 'size': Pt(36), 'bold': True},
                'subtitle': {'name': 'Segoe UI', 'size': Pt(24), 'bold': False},
                'content': {'name': 'Segoe UI', 'size': Pt(18), 'bold': False}
            },
            'marketing': {
                'title': {'name': 'Arial', 'size': Pt(40), 'bold': True},
                'subtitle': {'name': 'Arial', 'size': Pt(26), 'bold': False},
                'content': {'name': 'Arial', 'size': Pt(20), 'bold': False}
            },
            'strategy': {
                'title': {'name': 'Calibri', 'size': Pt(38), 'bold': True},
                'subtitle': {'name': 'Calibri', 'size': Pt(24), 'bold': False},
                'content': {'name': 'Calibri', 'size': Pt(18), 'bold': False}
            }
        }
        
        # Store configurations for later use
        self.current_colors = color_schemes.get(presentation_style, color_schemes['executive'])
        self.current_fonts = font_configs.get(presentation_style, font_configs['executive'])
        
        # Apply background styling if possible
        try:
            slide_master = ppt.slide_master
            background = slide_master.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = self.current_colors['background']
        except Exception as e:
            logger.warning(f"Could not apply master slide styling: {e}")
    
    def _create_slide(self, ppt: Presentation, slide_data: Dict[str, Any], presentation_style: str):
        """Create individual slide from parsed data."""
        
        slide_type = slide_data.get('type', 'content_slide')
        
        if slide_type == 'title_slide':
            self._create_title_slide(ppt, slide_data)
        elif slide_type == 'content_slide':
            self._create_content_slide(ppt, slide_data)
        elif slide_type == 'use_case_slide':
            self._create_use_case_slide(ppt, slide_data)
        else:
            self._create_content_slide(ppt, slide_data)  # Default to content slide
    
    def _create_title_slide(self, ppt: Presentation, slide_data: Dict[str, Any]):
        """Create title slide with proper template styling."""
        
        slide_layout = ppt.slide_layouts[0]  # Title slide layout
        slide = ppt.slides.add_slide(slide_layout)
        
        # Set title
        title = slide.shapes.title
        title.text = slide_data.get('title', 'Presentation Title')
        title.text_frame.paragraphs[0].font.size = self.current_fonts['title']['size']
        title.text_frame.paragraphs[0].font.color.rgb = self.current_colors['primary']
        title.text_frame.paragraphs[0].font.bold = self.current_fonts['title']['bold']
        title.text_frame.paragraphs[0].font.name = self.current_fonts['title']['name']
        
        # Set subtitle if available
        if len(slide.placeholders) > 1:
            subtitle = slide.placeholders[1]
            subtitle.text = slide_data.get('subtitle', 'Professional Business Presentation')
            subtitle.text_frame.paragraphs[0].font.size = self.current_fonts['subtitle']['size']
            subtitle.text_frame.paragraphs[0].font.color.rgb = self.current_colors['text']
            subtitle.text_frame.paragraphs[0].font.name = self.current_fonts['subtitle']['name']
    
    def _create_content_slide(self, ppt: Presentation, slide_data: Dict[str, Any]):
        """Create content slide with bullet points."""
        
        slide_layout = ppt.slide_layouts[1]  # Content slide layout
        slide = ppt.slides.add_slide(slide_layout)
        
        # Set title
        title = slide.shapes.title
        title.text = slide_data.get('title', 'Slide Title')
        title.text_frame.paragraphs[0].font.size = self.current_fonts['title']['size']
        title.text_frame.paragraphs[0].font.color.rgb = self.current_colors['primary']
        title.text_frame.paragraphs[0].font.bold = self.current_fonts['title']['bold']
        title.text_frame.paragraphs[0].font.name = self.current_fonts['title']['name']
        
        # Add bullet points
        if len(slide.placeholders) > 1:
            content = slide.placeholders[1]
            content.text_frame.clear()
            
            bullet_points = slide_data.get('bullet_points', [])
            content_blocks = slide_data.get('content_blocks', [])
            
            # Combine bullet points and content blocks
            all_content = bullet_points + content_blocks
            
            for i, point in enumerate(all_content[:5]):  # Max 5 points
                if i == 0:
                    p = content.text_frame.paragraphs[0]
                else:
                    p = content.text_frame.add_paragraph()
                
                # Clean bullet point text
                bullet_text = point.strip()
                if not bullet_text.startswith('•'):
                    bullet_text = f"• {bullet_text}"
                
                p.text = bullet_text
                p.font.size = self.current_fonts['content']['size']
                p.font.color.rgb = self.current_colors['text']
                p.font.name = self.current_fonts['content']['name']
                p.level = 0
                p.space_before = Pt(6)
                p.space_after = Pt(6)
        
        # Add notes
        if slide_data.get('notes'):
            slide.notes_slide.notes_text_frame.text = slide_data['notes']
    
    def _create_use_case_slide(self, ppt: Presentation, slide_data: Dict[str, Any]):
        """Create specialized use case slide."""
        
        slide_layout = ppt.slide_layouts[1]  # Content slide layout
        slide = ppt.slides.add_slide(slide_layout)
        
        # Set title
        title = slide.shapes.title
        title.text = slide_data.get('title', 'Use Case')
        title.text_frame.paragraphs[0].font.size = self.current_fonts['title']['size']
        title.text_frame.paragraphs[0].font.color.rgb = self.current_colors['primary']
        title.text_frame.paragraphs[0].font.bold = self.current_fonts['title']['bold']
        title.text_frame.paragraphs[0].font.name = self.current_fonts['title']['name']
        
        # Create structured content
        if len(slide.placeholders) > 1:
            content = slide.placeholders[1]
            content.text_frame.clear()
            
            # Add content blocks (Problem, Solution, Benefits)
            content_blocks = slide_data.get('content_blocks', [])
            bullet_points = slide_data.get('bullet_points', [])
            
            all_content = content_blocks + bullet_points
            
            for i, block in enumerate(all_content[:5]):
                if i == 0:
                    p = content.text_frame.paragraphs[0]
                else:
                    p = content.text_frame.add_paragraph()
                
                # Format content block
                block_text = block.strip()
                if not block_text.startswith('•') and not any(block_text.startswith(prefix) for prefix in ['Problem:', 'Solution:', 'Benefits:']):
                    block_text = f"• {block_text}"
                
                p.text = block_text
                p.font.size = self.current_fonts['content']['size']
                p.font.color.rgb = self.current_colors['text']
                p.font.name = self.current_fonts['content']['name']
                p.level = 0
                p.space_before = Pt(6)
                p.space_after = Pt(6)
        
        # Add notes
        if slide_data.get('notes'):
            slide.notes_slide.notes_text_frame.text = slide_data['notes']
    
    def _upload_pptx_to_s3(self, pptx_path: str, session_id: str, company_name: str, presentation_style: str) -> Optional[str]:
        """Upload PowerPoint to S3 and return URL."""
        
        try:
            timestamp = datetime.utcnow().strftime("%Y%m%d_%H%M%S")
            s3_key = f"presentations/{session_id}/{presentation_style}/{timestamp}_presentation.pptx"
            
            # Upload to S3
            s3_client.upload_file(
                pptx_path,
                S3_BUCKET,
                s3_key,
                ExtraArgs={
                    'ContentType': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
                    'Metadata': {
                        'session_id': session_id,
                        'company_name': company_name,
                        'presentation_style': presentation_style,
                        'generated_at': timestamp,
                        'editable': 'true'
                    }
                }
            )
            
            # Generate object URL
            s3_url = f"https://{S3_BUCKET}.s3.amazonaws.com/{s3_key}"
            
            logger.info(f"✅ PowerPoint uploaded to S3: {s3_key}")
            return s3_url
            
        except Exception as e:
            logger.error(f"❌ S3 upload failed: {e}")
            return None
    
    def _create_fallback_presentation_xml(self, company_profile: CompanyProfile, 
                                        use_cases: List[UseCaseStructured],
                                        research_data: Dict[str, Any], 
                                        presentation_style: str) -> str:
        """Create fallback XML presentation content."""
        
        fallback_xml = f"""
<presentation_title>Business Transformation Strategy for {company_profile.name}</presentation_title>
<template_type>{presentation_style}</template_type>

<slide type="title_slide">
<slide_title>Business Transformation Strategy for {company_profile.name}</slide_title>
<subtitle>Leveraging Technology for Competitive Advantage</subtitle>
<speaker_notes>Welcome and introduction to transformation opportunities</speaker_notes>
</slide>

<slide type="content_slide">
<slide_title>Executive Overview</slide_title>
<bullet>{company_profile.name} operates in the {company_profile.industry} sector</bullet>
<bullet>Current growth stage: {company_profile.growth_stage}</bullet>
<bullet>Technology maturity: {company_profile.cloud_maturity}</bullet>
<bullet>Identified {len(use_cases)} strategic transformation opportunities</bullet>
<speaker_notes>Brief company background and transformation readiness assessment</speaker_notes>
</slide>
        """
        
        # Add use case slides (minimum 3)
        for i, uc in enumerate(use_cases[:3], 1):
            fallback_xml += f"""
<slide type="use_case_slide">
<slide_title>Use Case {i}: {uc.title}</slide_title>
<content_block>Challenge: {uc.current_state}</content_block>
<content_block>Solution: {uc.proposed_solution[:150]}...</content_block>
<content_block>Benefits: {uc.business_value}</content_block>
<bullet>Timeline: {uc.timeline_months} months implementation</bullet>
<bullet>Investment: ${uc.monthly_cost_usd:,}/month operational cost</bullet>
<speaker_notes>Detailed explanation of {uc.title} implementation and benefits</speaker_notes>
</slide>
            """
        
        fallback_xml += """
<slide type="content_slide">
<slide_title>Implementation Approach</slide_title>
<bullet>Phase 1: Foundation and assessment (Months 1-2)</bullet>
<bullet>Phase 2: Core implementation (Months 3-6)</bullet>
<bullet>Phase 3: Optimization and scaling (Months 7-12)</bullet>
<bullet>Continuous improvement and value measurement</bullet>
<speaker_notes>Structured implementation methodology with risk mitigation</speaker_notes>
</slide>

<slide type="content_slide">
<slide_title>Next Steps</slide_title>
<bullet>Conduct detailed assessment of current capabilities</bullet>
<bullet>Prioritize use cases based on business impact</bullet>
<bullet>Develop detailed implementation roadmap</bullet>
<bullet>Establish success metrics and governance</bullet>
<speaker_notes>Action items and follow-up recommendations</speaker_notes>
</slide>
        """
        
        return fallback_xml