"""
Local-only PPT test without AWS dependencies.
"""
import sys
import os
from datetime import datetime

# Add project paths
project_root = os.path.dirname(__file__)
sys.path.insert(0, project_root)
sys.path.insert(0, os.path.join(project_root, 'src'))

# Mock AWS modules before importing
class MockS3Client:
    def upload_file(self, *args, **kwargs):
        print("Mock S3 upload - file saved locally")
        return True

class MockStatusTracker:
    def __init__(self, session_id):
        self.session_id = session_id
    
    def update_status(self, *args, **kwargs):
        pass  # Do nothing

# Mock AWS services
sys.modules['boto3'] = type('MockBoto3', (), {
    'client': lambda x, **kwargs: MockS3Client(),
    'resource': lambda x, **kwargs: None
})()

def test_direct_ppt_creation():
    """Test direct PPT creation without the full system."""
    
    print("Creating PowerPoint presentation directly...")
    
    try:
        from pptx import Presentation
        from pptx.util import Pt
        from pptx.dml.color import RGBColor
        
        # Create presentation
        ppt = Presentation()
        
        # Executive template colors
        navy_blue = RGBColor(30, 58, 138)
        dark_gray = RGBColor(31, 41, 55)
        
        # Title slide
        title_layout = ppt.slide_layouts[0]
        slide = ppt.slides.add_slide(title_layout)
        
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "Business Transformation Strategy"
        subtitle.text = "AI-Powered Solutions for Digital Excellence"
        
        # Style title
        title.text_frame.paragraphs[0].font.size = Pt(40)
        title.text_frame.paragraphs[0].font.color.rgb = navy_blue
        title.text_frame.paragraphs[0].font.bold = True
        
        # Style subtitle
        subtitle.text_frame.paragraphs[0].font.size = Pt(24)
        subtitle.text_frame.paragraphs[0].font.color.rgb = dark_gray
        
        # Content slide 1 - Overview
        content_layout = ppt.slide_layouts[1]
        slide = ppt.slides.add_slide(content_layout)
        
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "Executive Overview"
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.color.rgb = navy_blue
        title.text_frame.paragraphs[0].font.bold = True
        
        content.text = """• Strategic transformation through AI and automation
- Focus on operational efficiency and customer experience
- Identified 3 high-impact use cases for implementation
- Conservative ROI projections with realistic timelines
- Phased approach to minimize business disruption"""
        
        # Use Case 1
        slide = ppt.slides.add_slide(content_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "Use Case 1: Intelligent Process Automation"
        title.text_frame.paragraphs[0].font.size = Pt(32)
        title.text_frame.paragraphs[0].font.color.rgb = navy_blue
        
        content.text = """Challenge: Manual processes causing inefficiencies
Solution: AI-powered automation for key workflows
Benefits: 30-40% reduction in processing time
Timeline: 4-6 months implementation
Success Metrics: Efficiency gains and error reduction"""
        
        # Use Case 2
        slide = ppt.slides.add_slide(content_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "Use Case 2: Predictive Analytics Platform"
        title.text_frame.paragraphs[0].font.size = Pt(32)
        title.text_frame.paragraphs[0].font.color.rgb = navy_blue
        
        content.text = """Challenge: Reactive decision-making with limited insights
Solution: Advanced analytics for predictive intelligence
Benefits: 25% improvement in decision accuracy
Timeline: 6-8 months implementation
Success Metrics: Decision quality and response time"""
        
        # Use Case 3
        slide = ppt.slides.add_slide(content_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "Use Case 3: Customer Experience Enhancement"
        title.text_frame.paragraphs[0].font.size = Pt(32)
        title.text_frame.paragraphs[0].font.color.rgb = navy_blue
        
        content.text = """Challenge: Limited personalization and engagement
Solution: AI-driven customer experience optimization
Benefits: 20% increase in customer satisfaction
Timeline: 5-7 months implementation
Success Metrics: NPS scores and retention rates"""
        
        # Implementation slide
        slide = ppt.slides.add_slide(content_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "Implementation Roadmap"
        title.text_frame.paragraphs[0].font.size = Pt(32)
        title.text_frame.paragraphs[0].font.color.rgb = navy_blue
        
        content.text = """• Phase 1: Assessment and Planning (Months 1-2)
- Phase 2: Core Implementation (Months 3-8)
- Phase 3: Testing and Optimization (Months 9-10)
- Phase 4: Full Deployment (Months 11-12)
- Ongoing: Monitoring and Enhancement"""
        
        # Next steps slide
        slide = ppt.slides.add_slide(content_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "Next Steps & Recommendations"
        title.text_frame.paragraphs[0].font.size = Pt(32)
        title.text_frame.paragraphs[0].font.color.rgb = navy_blue
        
        content.text = """• Approve transformation budget and resources
- Establish project governance and team structure
- Begin detailed assessment of current capabilities
- Select technology partners and vendors
- Communicate transformation vision to stakeholders"""
        
        # Save locally
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        filename = f"business_transformation_{timestamp}.pptx"
        ppt.save(filename)
        
        print(f"SUCCESS! PowerPoint created: {filename}")
        print(f"Location: {os.path.abspath(filename)}")
        print(f"Slides: {len(ppt.slides)}")
        print(f"Template: Executive (Navy blue theme)")
        print("\nOpen this file in PowerPoint to view and edit!")
        
        return filename
        
    except Exception as e:
        print(f"Error: {e}")
        import traceback
        traceback.print_exc()

if __name__ == "__main__":
    test_direct_ppt_creation()