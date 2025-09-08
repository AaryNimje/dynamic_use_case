#!/usr/bin/env python3
"""
Create template PowerPoint files for each presentation type
This generates empty template .pptx files with the proper styling and structure
"""
import os
from datetime import datetime

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    PPT_AVAILABLE = True
except ImportError:
    PPT_AVAILABLE = False

def create_template_directory():
    """Create templates directory if it doesn't exist"""
    template_dir = "templates"
    if not os.path.exists(template_dir):
        os.makedirs(template_dir)
        print(f"Created templates directory: {template_dir}")
    return template_dir

def create_first_deck_template():
    """Create First Deck Call template"""
    ppt = Presentation()
    
    # Color scheme - Deep navy, professional
    primary = RGBColor(20, 33, 61)     # Deep navy
    secondary = RGBColor(52, 73, 94)   # Slate blue
    accent = RGBColor(230, 126, 34)    # Orange
    text = RGBColor(44, 62, 80)        # Dark blue-gray
    
    # Title slide
    slide = ppt.slides.add_slide(ppt.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Strategic Partnership Opportunity"
    subtitle.text = "[Company Name] Executive Overview"
    
    # Style title slide
    title.text_frame.paragraphs[0].font.size = Pt(48)
    title.text_frame.paragraphs[0].font.color.rgb = primary
    title.text_frame.paragraphs[0].font.bold = True
    
    subtitle.text_frame.paragraphs[0].font.size = Pt(28)
    subtitle.text_frame.paragraphs[0].font.color.rgb = secondary
    
    # Company overview slide
    slide = ppt.slides.add_slide(ppt.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Company Overview"
    content.text = "• Business profile and market position\n• Key metrics and performance indicators\n• Current challenges and opportunities\n• Strategic objectives and priorities"
    
    # Style content slide
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.color.rgb = primary
    title.text_frame.paragraphs[0].font.bold = True
    
    for paragraph in content.text_frame.paragraphs:
        paragraph.font.size = Pt(24)
        paragraph.font.color.rgb = text
        paragraph.space_before = Pt(8)
    
    # Strategic opportunities slide
    slide = ppt.slides.add_slide(ppt.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Strategic Opportunities Identified"
    content.text = "• Primary transformation opportunity with business impact\n• Secondary optimization area for competitive advantage\n• Innovation initiative for market leadership\n• Partnership potential for accelerated growth"
    
    # Style
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.color.rgb = primary
    title.text_frame.paragraphs[0].font.bold = True
    
    for paragraph in content.text_frame.paragraphs:
        paragraph.font.size = Pt(24)
        paragraph.font.color.rgb = text
        paragraph.space_before = Pt(8)
    
    # Value proposition slide
    slide = ppt.slides.add_slide(ppt.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Potential Business Value"
    content.text = "• Cost optimization potential: 15-25% efficiency gains\n• Revenue growth enablement through innovation\n• Competitive advantage through strategic positioning\n• Partnership value creation and market expansion"
    
    # Style value slide
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.color.rgb = primary
    title.text_frame.paragraphs[0].font.bold = True
    
    for paragraph in content.text_frame.paragraphs:
        paragraph.font.size = Pt(24)
        paragraph.font.color.rgb = text
        paragraph.space_before = Pt(8)
    
    # Next steps slide
    slide = ppt.slides.add_slide(ppt.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Proposed Next Steps"
    content.text = "• Detailed assessment and discovery phase\n• Proof of concept development\n• Implementation roadmap creation\n• Strategic partnership discussion"
    
    # Style next steps slide
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.color.rgb = primary
    title.text_frame.paragraphs[0].font.bold = True
    
    for paragraph in content.text_frame.paragraphs:
        paragraph.font.size = Pt(24)
        paragraph.font.color.rgb = text
        paragraph.space_before = Pt(8)
    
    return ppt

def create_marketing_template():
    """Create Marketing presentation template"""
    ppt = Presentation()
    
    # Color scheme - Vibrant, engaging
    primary = RGBColor(225, 45, 139)    # Vibrant pink
    secondary = RGBColor(74, 144, 226)  # Bright blue
    accent = RGBColor(255, 193, 7)      # Golden yellow
    text = RGBColor(33, 37, 41)         # Dark gray
    success = RGBColor(40, 167, 69)     # Success green
    
    # Title slide
    slide = ppt.slides.add_slide(ppt.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Transform Your Business Today"
    subtitle.text = "Unlock Growth and Innovation with [Solution Name]"
    
    # Style title slide
    title.text_frame.paragraphs[0].font.size = Pt(52)
    title.text_frame.paragraphs[0].font.color.rgb = primary
    title.text_frame.paragraphs[0].font.bold = True
    
    subtitle.text_frame.paragraphs[0].font.size = Pt(30)
    subtitle.text_frame.paragraphs[0].font.color.rgb = secondary
    subtitle.text_frame.paragraphs[0].font.bold = True
    
    # Problem slide
    slide = ppt.slides.add_slide(ppt.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "The Challenge Every Business Faces"
    content.text = "• Rising costs and competitive pressure\n• Inefficient processes holding you back\n• Missing growth opportunities\n• Technology gaps limiting potential\n\n⚠️ The cost of inaction grows every day"
    
    # Style problem slide
    title.text_frame.paragraphs[0].font.size = Pt(42)
    title.text_frame.paragraphs[0].font.color.rgb = primary
    title.text_frame.paragraphs[0].font.bold = True
    
    for paragraph in content.text_frame.paragraphs:
        paragraph.font.size = Pt(22)
        paragraph.font.color.rgb = accent
        paragraph.space_before = Pt(10)
        paragraph.font.bold = True
    
    # Solution slide
    slide = ppt.slides.add_slide(ppt.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "The Transformation Solution"
    content.text = "• Streamline operations for maximum efficiency\n• Unlock hidden revenue opportunities\n• Leverage cutting-edge technology\n• Gain competitive market advantage\n\n✨ Transform challenges into success stories"
    
    # Style solution slide
    title.text_frame.paragraphs[0].font.size = Pt(42)
    title.text_frame.paragraphs[0].font.color.rgb = primary
    title.text_frame.paragraphs[0].font.bold = True
    
    for paragraph in content.text_frame.paragraphs:
        paragraph.font.size = Pt(22)
        paragraph.font.color.rgb = success
        paragraph.space_before = Pt(10)
        paragraph.font.bold = True
    
    # Benefits slide
    slide = ppt.slides.add_slide(ppt.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Real Results You Can Achieve"
    content.text = "• Operational Excellence: 30% efficiency improvement\n• Financial Impact: 20% cost reduction potential\n• Market Position: Competitive advantage gains\n• Innovation Leadership: Technology-driven growth\n\n🚀 Start seeing results in 90 days"
    
    # Style benefits slide
    title.text_frame.paragraphs[0].font.size = Pt(42)
    title.text_frame.paragraphs[0].font.color.rgb = primary
    title.text_frame.paragraphs[0].font.bold = True
    
    for paragraph in content.text_frame.paragraphs:
        paragraph.font.size = Pt(22)
        paragraph.font.color.rgb = success
        paragraph.space_before = Pt(10)
        paragraph.font.bold = True
    
    # Success stories slide
    slide = ppt.slides.add_slide(ppt.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Success Stories and Use Cases"
    content.text = "• Manufacturing: 40% production efficiency gains\n• Healthcare: 50% patient processing improvement\n• Finance: 60% compliance automation success\n• Retail: 35% customer satisfaction increase\n\n📈 Join the transformation leaders"
    
    # Style success slide
    title.text_frame.paragraphs[0].font.size = Pt(42)
    title.text_frame.paragraphs[0].font.color.rgb = primary
    title.text_frame.paragraphs[0].font.bold = True
    
    for paragraph in content.text_frame.paragraphs:
        paragraph.font.size = Pt(22)
        paragraph.font.color.rgb = text
        paragraph.space_before = Pt(10)
    
    # Call to action slide
    slide = ppt.slides.add_slide(ppt.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Start Your Transformation Journey"
    content.text = "• Schedule your free transformation assessment\n• Join our pilot program for early adopters\n• Explore partnership opportunities\n• Get exclusive access to innovation previews\n\n🎯 Contact us today to begin!"
    
    # Style CTA slide
    title.text_frame.paragraphs[0].font.size = Pt(42)
    title.text_frame.paragraphs[0].font.color.rgb = primary
    title.text_frame.paragraphs[0].font.bold = True
    
    for paragraph in content.text_frame.paragraphs:
        paragraph.font.size = Pt(22)
        paragraph.font.color.rgb = accent
        paragraph.space_before = Pt(10)
        paragraph.font.bold = True
    
    return ppt

def create_use_case_template():
    """Create Use Case presentation template"""
    ppt = Presentation()
    
    # Color scheme - Professional, detailed
    primary = RGBColor(99, 102, 241)    # Indigo
    secondary = RGBColor(139, 69, 19)   # Saddle brown
    accent = RGBColor(245, 158, 11)     # Amber
    text = RGBColor(55, 65, 81)         # Gray-700
    
    # Title slide
    slide = ppt.slides.add_slide(ppt.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Use Case Implementation Strategy"
    subtitle.text = "[Company Name] Transformation Scenarios"
    
    # Style title slide
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = primary
    title.text_frame.paragraphs[0].font.bold = True
    
    subtitle.text_frame.paragraphs[0].font.size = Pt(26)
    subtitle.text_frame.paragraphs[0].font.color.rgb = secondary
    
    # Use case overview slide
    slide = ppt.slides.add_slide(ppt.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Use Case Portfolio Overview"
    content.text = "• Use Case 1: Process automation and efficiency optimization\n• Use Case 2: Data analytics and business intelligence\n• Use Case 3: Customer experience enhancement platform\n• Use Case 4: Supply chain and operations optimization"
    
    # Style content slide
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.color.rgb = primary
    title.text_frame.paragraphs[0].font.bold = True
    
    for paragraph in content.text_frame.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = text
        paragraph.space_before = Pt(6)
    
    # Detailed use case 1 slide
    slide = ppt.slides.add_slide(ppt.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Use Case 1: Process Automation"
    content.text = "Current State: Manual processes causing delays and errors\n\nSolution: Automated workflow and task management system\n\nImplementation:\n• Phase 1: Assessment and planning (months 1-2)\n• Phase 2: Pilot implementation (months 3-4)\n• Phase 3: Full rollout (months 5-8)\n\nExpected Benefits:\n• 40% reduction in processing time\n• 60% decrease in manual errors\n• $200K annual cost savings"
    
    # Style detailed slide
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.color.rgb = primary
    title.text_frame.paragraphs[0].font.bold = True
    
    for paragraph in content.text_frame.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = text
        paragraph.space_before = Pt(6)
    
    # Detailed use case 2 slide
    slide = ppt.slides.add_slide(ppt.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Use Case 2: Data Analytics Platform"
    content.text = "Current State: Limited visibility into business performance\n\nSolution: Real-time analytics and business intelligence dashboard\n\nImplementation:\n• Phase 1: Data integration and cleansing (months 1-3)\n• Phase 2: Dashboard development (months 4-5)\n• Phase 3: User training and adoption (months 6-7)\n\nExpected Benefits:\n• 50% faster decision-making\n• 25% improvement in forecast accuracy\n• $150K annual efficiency gains"
    
    # Style detailed slide 2
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.color.rgb = primary
    title.text_frame.paragraphs[0].font.bold = True
    
    for paragraph in content.text_frame.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = text
        paragraph.space_before = Pt(6)
    
    # Implementation methodology slide
    slide = ppt.slides.add_slide(ppt.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Implementation Methodology"
    content.text = "1. Discovery and requirements gathering\n2. Solution design and architecture\n3. Pilot development and testing\n4. Training and change management\n5. Full deployment and optimization\n\nTimeline: 12-18 month implementation cycle\nSuccess Metrics: ROI measurement and KPI tracking"
    
    # Style methodology slide
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.color.rgb = primary
    title.text_frame.paragraphs[0].font.bold = True
    
    for paragraph in content.text_frame.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = text
        paragraph.space_before = Pt(6)
    
    # Risk management slide
    slide = ppt.slides.add_slide(ppt.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Risk Management and Success Factors"
    content.text = "Risk Factors:\n• Technical integration challenges\n• User adoption and change resistance\n• Data migration and quality issues\n\nMitigation Strategies:\n• Phased implementation approach\n• Comprehensive training program\n• Dedicated project management office\n\nSuccess Factors:\n• Executive sponsorship and support\n• Cross-functional team collaboration\n• Continuous monitoring and optimization"
    
    # Style risk slide
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.color.rgb = primary
    title.text_frame.paragraphs[0].font.bold = True
    
    for paragraph in content.text_frame.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = text
        paragraph.space_before = Pt(6)
    
    return ppt

def create_technical_template():
    """Create Technical Architecture presentation template"""
    ppt = Presentation()
    
    # Color scheme - Clean, technical
    primary = RGBColor(30, 41, 59)      # Slate-800
    secondary = RGBColor(71, 85, 105)   # Slate-600
    accent = RGBColor(14, 165, 233)     # Sky-500
    text = RGBColor(51, 65, 85)         # Slate-700
    code = RGBColor(239, 68, 68)        # Red-500
    
    # Title slide
    slide = ppt.slides.add_slide(ppt.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Technical Architecture Overview"
    subtitle.text = "[Company Name] System Design and Implementation"
    
    # Style title slide
    title.text_frame.paragraphs[0].font.size = Pt(42)
    title.text_frame.paragraphs[0].font.color.rgb = primary
    title.text_frame.paragraphs[0].font.bold = True
    
    subtitle.text_frame.paragraphs[0].font.size = Pt(24)
    subtitle.text_frame.paragraphs[0].font.color.rgb = secondary
    
    # Architecture slide
    slide = ppt.slides.add_slide(ppt.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "System Architecture Components"
    content.text = "Architecture Layers:\n• Frontend: React/Vue.js with responsive design\n• Backend: Node.js/Python microservices architecture\n• Database: PostgreSQL with Redis caching layer\n• Cloud: AWS/Azure with container orchestration\n\nDesign Principles:\n• Scalable and modular architecture\n• API-first development approach\n• Cloud-native deployment strategy"
    
    # Style architecture slide
    title.text_frame.paragraphs[0].font.size = Pt(34)
    title.text_frame.paragraphs[0].font.color.rgb = primary
    title.text_frame.paragraphs[0].font.bold = True
    
    for paragraph in content.text_frame.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = text
        paragraph.space_before = Pt(4)
    
    # Technology stack slide
    slide = ppt.slides.add_slide(ppt.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Technology Stack and Platform"
    content.text = "Frontend Technologies:\n• React 18 with TypeScript\n• Tailwind CSS for styling\n• Redux for state management\n\nBackend Technologies:\n• Node.js with Express framework\n• GraphQL API layer\n• JWT authentication\n\nInfrastructure:\n• Docker containerization\n• Kubernetes orchestration\n• CI/CD with GitHub Actions"
    
    # Style tech stack slide
    title.text_frame.paragraphs[0].font.size = Pt(34)
    title.text_frame.paragraphs[0].font.color.rgb = primary
    title.text_frame.paragraphs[0].font.bold = True
    
    for paragraph in content.text_frame.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = text
        paragraph.space_before = Pt(4)
    
    # Integration slide
    slide = ppt.slides.add_slide(ppt.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Integration Architecture"
    content.text = "Integration Patterns:\n• API-first design approach\n• Event-driven architecture\n• Real-time data synchronization\n\nData Flow:\n• Data ingestion and validation\n• Processing and transformation\n• Storage and retrieval optimization\n\nSecurity Measures:\n• End-to-end encryption\n• Multi-factor authentication\n• Regular security audits"
    
    # Style integration slide
    title.text_frame.paragraphs[0].font.size = Pt(34)
    title.text_frame.paragraphs[0].font.color.rgb = primary
    title.text_frame.paragraphs[0].font.bold = True
    
    for paragraph in content.text_frame.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = text
        paragraph.space_before = Pt(4)
    
    # Performance and security slide
    slide = ppt.slides.add_slide(ppt.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Performance and Security Requirements"
    content.text = "Performance Requirements:\n• Response time: <200ms for API calls\n• Throughput: 10,000+ concurrent users\n• Availability: 99.9% uptime SLA\n• Scalability: Auto-scaling based on demand\n\nSecurity Framework:\n• OWASP compliance standards\n• Data encryption at rest and in transit\n• Regular penetration testing\n• Compliance with GDPR and SOC 2"
    
    # Style performance slide
    title.text_frame.paragraphs[0].font.size = Pt(34)
    title.text_frame.paragraphs[0].font.color.rgb = primary
    title.text_frame.paragraphs[0].font.bold = True
    
    for paragraph in content.text_frame.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = text
        paragraph.space_before = Pt(4)
    
    return ppt

def create_strategy_template():
    """Create Strategy Planning presentation template"""
    ppt = Presentation()
    
    # Color scheme - Strategic, authoritative
    primary = RGBColor(79, 70, 229)     # Indigo-600
    secondary = RGBColor(107, 114, 128) # Gray-500
    accent = RGBColor(16, 185, 129)     # Emerald-500
    text = RGBColor(17, 24, 39)         # Gray-900
    highlight = RGBColor(245, 101, 101) # Red-400
    
    # Title slide
    slide = ppt.slides.add_slide(ppt.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Strategic Transformation Roadmap"
    subtitle.text = "[Company Name] 3-Year Strategic Plan"
    
    # Style title slide
    title.text_frame.paragraphs[0].font.size = Pt(46)
    title.text_frame.paragraphs[0].font.color.rgb = primary
    title.text_frame.paragraphs[0].font.bold = True
    
    subtitle.text_frame.paragraphs[0].font.size = Pt(28)
    subtitle.text_frame.paragraphs[0].font.color.rgb = secondary
    
    # Current state slide
    slide = ppt.slides.add_slide(ppt.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Current State Assessment"
    content.text = "Organizational Strengths:\n• Market leadership in core segments\n• Strong customer relationships and brand recognition\n• Experienced leadership team and workforce\n\nStrategic Challenges:\n• Technology infrastructure modernization needs\n• Operational efficiency optimization opportunities\n• Market expansion and competitive positioning"
    
    # Style current state slide
    title.text_frame.paragraphs[0].font.size = Pt(38)
    title.text_frame.paragraphs[0].font.color.rgb = primary
    title.text_frame.paragraphs[0].font.bold = True
    
    for paragraph in content.text_frame.paragraphs:
        paragraph.font.size = Pt(20)
        paragraph.font.color.rgb = text
        paragraph.space_before = Pt(7)
    
    # Strategic vision slide
    slide = ppt.slides.add_slide(ppt.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Strategic Vision and Objectives"
    content.text = "Vision Statement:\nTo become the industry leader in innovation and customer satisfaction through strategic transformation and operational excellence.\n\nStrategic Objectives:\n• Achieve 25% market share growth over 3 years\n• Implement digital transformation across all operations\n• Establish new revenue streams and partnerships\n• Build sustainable competitive advantages"
    
    # Style vision slide
    title.text_frame.paragraphs[0].font.size = Pt(38)
    title.text_frame.paragraphs[0].font.color.rgb = primary
    title.text_frame.paragraphs[0].font.bold = True
    
    for paragraph in content.text_frame.paragraphs:
        paragraph.font.size = Pt(20)
        paragraph.font.color.rgb = text
        paragraph.space_before = Pt(7)
    
    # Strategic initiatives slide
    slide = ppt.slides.add_slide(ppt.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Strategic Initiatives Portfolio"
    content.text = "Growth Initiatives:\n• Market expansion into emerging regions\n• New product development and innovation\n• Strategic partnerships and acquisitions\n\nEfficiency Initiatives:\n• Process automation and optimization\n• Technology infrastructure modernization\n• Organizational restructuring for agility\n\nInnovation Initiatives:\n• Digital transformation programs\n• Research and development investments\n• Customer experience enhancement"
    
    # Style initiatives slide
    title.text_frame.paragraphs[0].font.size = Pt(38)
    title.text_frame.paragraphs[0].font.color.rgb = primary
    title.text_frame.paragraphs[0].font.bold = True
    
    for paragraph in content.text_frame.paragraphs:
        paragraph.font.size = Pt(20)
        paragraph.font.color.rgb = text
        paragraph.space_before = Pt(7)
    
    # Implementation roadmap slide
    slide = ppt.slides.add_slide(ppt.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "3-Year Implementation Roadmap"
    content.text = "Year 1 - Foundation Building:\n• Infrastructure modernization\n• Process optimization initiatives\n• Team capability development\n\nYear 2 - Growth Acceleration:\n• Market expansion execution\n• New product launches\n• Partnership establishment\n\nYear 3 - Market Leadership:\n• Innovation commercialization\n• Competitive advantage realization\n• Sustainable growth achievement"
    
    # Style roadmap slide
    title.text_frame.paragraphs[0].font.size = Pt(38)
    title.text_frame.paragraphs[0].font.color.rgb = primary
    title.text_frame.paragraphs[0].font.bold = True
    
    for paragraph in content.text_frame.paragraphs:
        paragraph.font.size = Pt(20)
        paragraph.font.color.rgb = text
        paragraph.space_before = Pt(7)
    
    # Success metrics slide
    slide = ppt.slides.add_slide(ppt.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Success Metrics and Governance"
    content.text = "Financial Performance Indicators:\n• Revenue growth: 25% CAGR target\n• Profit margin improvement: 15% increase\n• ROI on strategic investments: >20%\n\nOperational Excellence Metrics:\n• Customer satisfaction: >90% target\n• Process efficiency: 30% improvement\n• Employee engagement: Top quartile\n\nGovernance Structure:\n• Monthly strategic review meetings\n• Quarterly board reporting\n• Annual strategic plan updates"
    
    # Style metrics slide
    title.text_frame.paragraphs[0].font.size = Pt(38)
    title.text_frame.paragraphs[0].font.color.rgb = primary
    title.text_frame.paragraphs[0].font.bold = True
    
    for paragraph in content.text_frame.paragraphs:
        paragraph.font.size = Pt(20)
        paragraph.font.color.rgb = text
        paragraph.space_before = Pt(7)
    
    return ppt

def save_template(ppt, template_name, template_dir):
    """Save template to file"""
    filename = f"{template_name}_template.pptx"
    filepath = os.path.join(template_dir, filename)
    ppt.save(filepath)
    
    file_size = os.path.getsize(filepath)
    print(f"Created {template_name} template: {filename} ({file_size:,} bytes)")
    
    return filepath

def main():
    """Create all template files"""
    
    print("Creating PowerPoint Template Files")
    print("=" * 35)
    
    if not PPT_AVAILABLE:
        print("Missing PowerPoint library. Install with:")
        print("pip install python-pptx")
        return
    
    # Create templates directory
    template_dir = create_template_directory()
    
    # Create all templates
    templates = [
        ("first_deck", create_first_deck_template),
        ("marketing", create_marketing_template),
        ("use_case", create_use_case_template),
        ("technical", create_technical_template),
        ("strategy", create_strategy_template)
    ]
    
    created_files = []
    
    for template_name, create_func in templates:
        try:
            print(f"\nCreating {template_name} template...")
            ppt = create_func()
            filepath = save_template(ppt, template_name, template_dir)
            created_files.append(filepath)
        except Exception as e:
            print(f"Error creating {template_name} template: {e}")
    
    print(f"\n" + "=" * 35)
    print(f"Template Creation Complete!")
    print(f"Created {len(created_files)} template files in '{template_dir}' folder:")
    
    for filepath in created_files:
        print(f"  • {os.path.basename(filepath)}")
    
    print(f"\nThese templates can be used as:")
    print(f"1. Starting points for manual presentation creation")
    print(f"2. Reference for styling and structure")
    print(f"3. Base templates for the automated generation system")
    
    print(f"\nTemplate Details:")
    print(f"• First Deck: Executive overview (5 slides)")
    print(f"• Marketing: Persuasive presentation (6 slides)")
    print(f"• Use Case: Detailed scenarios (7 slides)")
    print(f"• Technical: Architecture specs (5 slides)")
    print(f"• Strategy: Strategic planning (6 slides)")

if __name__ == "__main__":
    main()