#!/usr/bin/env python3
"""
Generate PowerPoint presentations from PDF content using AWS Bedrock
Takes PDF content and creates a professional PPT with use cases and business insights
"""
import os
import boto3
import json
from datetime import datetime
from typing import Dict, Any, List, Optional

# PDF reading
try:
    import PyPDF2
    import fitz  # PyMuPDF
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

# PowerPoint creation
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    PPT_AVAILABLE = True
except ImportError:
    PPT_AVAILABLE = False

class PDFContentExtractor:
    """Extract text content from PDF files for analysis"""
    
    @staticmethod
    def extract_text(pdf_path: str) -> str:
        """Extract all text content from PDF"""
        
        if not PDF_AVAILABLE:
            raise ImportError("Install PDF libraries: pip install PyMuPDF PyPDF2")
        
        text_content = ""
        
        try:
            # Use PyMuPDF (better for most PDFs)
            doc = fitz.open(pdf_path)
            for page in doc:
                text_content += page.get_text()
            doc.close()
            
            if len(text_content.strip()) > 100:
                return PDFContentExtractor._clean_text(text_content)
        
        except Exception:
            pass
        
        try:
            # Fallback to PyPDF2
            with open(pdf_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page in pdf_reader.pages:
                    text_content += page.extract_text()
                    
            return PDFContentExtractor._clean_text(text_content)
            
        except Exception as e:
            raise Exception(f"Could not extract text from PDF: {e}")
    
    @staticmethod
    def _clean_text(text: str) -> str:
        """Clean extracted text"""
        import re
        # Remove excessive whitespace
        text = re.sub(r'\n\s*\n', '\n\n', text)
        text = re.sub(r' +', ' ', text)
        return text.strip()

class BedrockContentAnalyzer:
    """Analyze PDF content using AWS Bedrock to create PPT structure"""
    
    def __init__(self, region='us-east-1'):
        self.bedrock = boto3.client('bedrock-runtime', region_name=region)
        self.model_id = 'us.anthropic.claude-sonnet-4-20250514-v1:0'
    
    def analyze_content_for_ppt(self, content: str, company_name: str = None) -> Dict[str, Any]:
        """Analyze content and structure it for PowerPoint creation"""
        
        # Limit content to avoid token limits
        content_preview = content[:12000] if len(content) > 12000 else content
        
        analysis_prompt = f"""
You are a business presentation expert. Analyze this business document content and create a professional PowerPoint presentation structure.

DOCUMENT CONTENT TO ANALYZE:
{content_preview}

Your task: Extract key business insights and create an executive-ready presentation structure.

REQUIREMENTS:
1. Identify the main business focus and key transformation opportunities
2. Extract 3-5 specific use cases or strategic initiatives
3. Create slide content with clear titles and bullet points
4. Include realistic timelines and benefits (be conservative)
5. Structure for 8-12 professional slides

OUTPUT AS JSON:
{{
    "presentation_info": {{
        "title": "Professional presentation title based on content",
        "company": "{company_name or 'extracted from content'}",
        "focus_area": "main business area identified",
        "total_slides": 10
    }},
    "slides": [
        {{
            "slide_number": 1,
            "type": "title",
            "title": "Main Presentation Title",
            "subtitle": "Key focus area and date"
        }},
        {{
            "slide_number": 2, 
            "type": "overview",
            "title": "Executive Overview",
            "content": [
                "Key business insight 1",
                "Key business insight 2", 
                "Key business insight 3",
                "Strategic focus areas identified"
            ]
        }},
        {{
            "slide_number": 3,
            "type": "use_case",
            "title": "Use Case 1: [Specific Solution Name]",
            "challenge": "Current business challenge identified",
            "solution": "Proposed technology/process solution",
            "benefits": [
                "Specific benefit with realistic timeline",
                "Cost impact or efficiency gain",
                "Implementation consideration"
            ],
            "timeline": "6-12 months implementation"
        }}
        // Continue with more use case slides, implementation roadmap, and conclusion
    ]
}}

Focus on:
- Actionable business recommendations
- Realistic ROI estimates and timelines
- Professional executive-level content
- Clear problem-solution-benefit structure for use cases

Generate a complete presentation structure based on the actual content provided.
"""

        try:
            response = self.bedrock.invoke_model(
                modelId=self.model_id,
                body=json.dumps({
                    "anthropic_version": "bedrock-2023-05-31",
                    "max_tokens": 4000,
                    "messages": [{"role": "user", "content": analysis_prompt}]
                })
            )
            
            response_body = json.loads(response['body'].read())
            ai_response = response_body['content'][0]['text']
            
            # Extract JSON from response
            import re
            json_match = re.search(r'\{.*\}', ai_response, re.DOTALL)
            if json_match:
                return json.loads(json_match.group())
            else:
                return self._create_basic_structure(content, company_name)
                
        except Exception as e:
            print(f"Bedrock analysis failed: {e}")
            return self._create_basic_structure(content, company_name)
    
    def _create_basic_structure(self, content: str, company_name: str) -> Dict[str, Any]:
        """Create basic structure if Bedrock fails"""
        return {
            "presentation_info": {
                "title": f"{company_name or 'Business'} Strategic Analysis",
                "company": company_name or "Organization",
                "focus_area": "Business transformation and optimization",
                "total_slides": 8
            },
            "slides": [
                {
                    "slide_number": 1,
                    "type": "title",
                    "title": f"{company_name or 'Business'} Strategic Analysis",
                    "subtitle": f"Key Insights and Recommendations - {datetime.now().strftime('%B %Y')}"
                },
                {
                    "slide_number": 2,
                    "type": "overview", 
                    "title": "Executive Summary",
                    "content": [
                        "Comprehensive analysis of business operations completed",
                        "Key transformation opportunities identified",
                        "Strategic recommendations developed",
                        "Implementation roadmap outlined"
                    ]
                }
            ]
        }

class PowerPointGenerator:
    """Generate professional PowerPoint presentations"""
    
    def __init__(self):
        if not PPT_AVAILABLE:
            raise ImportError("Install PowerPoint library: pip install python-pptx")
        
        # Template color schemes
        self.templates = {
            "executive": {
                "primary": RGBColor(31, 78, 121),    # Navy blue
                "secondary": RGBColor(79, 129, 189), # Light blue
                "accent": RGBColor(192, 80, 77),     # Red
                "text": RGBColor(64, 64, 64)
            },
            "technical": {
                "primary": RGBColor(68, 84, 106),
                "secondary": RGBColor(112, 128, 144),
                "accent": RGBColor(70, 130, 180),
                "text": RGBColor(64, 64, 64)
            },
            "marketing": {
                "primary": RGBColor(46, 125, 50),
                "secondary": RGBColor(76, 175, 80),
                "accent": RGBColor(255, 193, 7),
                "text": RGBColor(64, 64, 64)
            }
        }
    
    def create_presentation(self, structure: Dict[str, Any], template: str = "executive") -> str:
        """Create PowerPoint from analyzed structure"""
        
        # Create new presentation
        ppt = Presentation()
        colors = self.templates.get(template, self.templates["executive"])
        
        # Create slides from structure
        for slide_data in structure.get('slides', []):
            self._create_slide(ppt, slide_data, colors)
        
        # Save presentation
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        company = structure.get('presentation_info', {}).get('company', 'business')
        company_clean = company.replace(' ', '_').lower()
        filename = f"{company_clean}_presentation_{timestamp}.pptx"
        
        ppt.save(filename)
        return filename
    
    def _create_slide(self, ppt: Presentation, slide_data: Dict[str, Any], colors: Dict[str, Any]):
        """Create individual slide"""
        
        slide_type = slide_data.get('type', 'content')
        
        if slide_type == 'title':
            self._create_title_slide(ppt, slide_data, colors)
        elif slide_type == 'overview':
            self._create_overview_slide(ppt, slide_data, colors)
        elif slide_type == 'use_case':
            self._create_use_case_slide(ppt, slide_data, colors)
        else:
            self._create_content_slide(ppt, slide_data, colors)
    
    def _create_title_slide(self, ppt: Presentation, data: Dict[str, Any], colors: Dict[str, Any]):
        """Create title slide"""
        slide_layout = ppt.slide_layouts[0]
        slide = ppt.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = data.get('title', 'Presentation Title')
        subtitle.text = data.get('subtitle', 'Subtitle')
        
        # Style title
        title.text_frame.paragraphs[0].font.size = Pt(44)
        title.text_frame.paragraphs[0].font.color.rgb = colors["primary"]
        title.text_frame.paragraphs[0].font.bold = True
        
        # Style subtitle
        subtitle.text_frame.paragraphs[0].font.size = Pt(24)
        subtitle.text_frame.paragraphs[0].font.color.rgb = colors["secondary"]
    
    def _create_overview_slide(self, ppt: Presentation, data: Dict[str, Any], colors: Dict[str, Any]):
        """Create overview slide with bullet points"""
        slide_layout = ppt.slide_layouts[1]
        slide = ppt.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = data.get('title', 'Overview')
        
        # Style title
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.color.rgb = colors["primary"]
        title.text_frame.paragraphs[0].font.bold = True
        
        # Add content
        content_items = data.get('content', [])
        if content_items:
            content.text = content_items[0]
            for item in content_items[1:]:
                p = content.text_frame.add_paragraph()
                p.text = item
                p.level = 0
        
        # Style content
        for paragraph in content.text_frame.paragraphs:
            paragraph.font.size = Pt(20)
            paragraph.font.color.rgb = colors["text"]
            paragraph.space_before = Pt(6)
    
    def _create_use_case_slide(self, ppt: Presentation, data: Dict[str, Any], colors: Dict[str, Any]):
        """Create use case slide"""
        slide_layout = ppt.slide_layouts[1]
        slide = ppt.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = data.get('title', 'Use Case')
        
        # Style title
        title.text_frame.paragraphs[0].font.size = Pt(32)
        title.text_frame.paragraphs[0].font.color.rgb = colors["primary"]
        title.text_frame.paragraphs[0].font.bold = True
        
        # Build content
        slide_content = []
        
        if data.get('challenge'):
            slide_content.append(f"Challenge: {data['challenge']}")
        
        if data.get('solution'):
            slide_content.append(f"Solution: {data['solution']}")
        
        if data.get('benefits'):
            slide_content.append("Benefits:")
            for benefit in data['benefits']:
                slide_content.append(f"  • {benefit}")
        
        if data.get('timeline'):
            slide_content.append(f"Timeline: {data['timeline']}")
        
        # Add to slide
        if slide_content:
            content.text = slide_content[0]
            for item in slide_content[1:]:
                p = content.text_frame.add_paragraph()
                p.text = item
                p.level = 0
        
        # Style content
        for paragraph in content.text_frame.paragraphs:
            paragraph.font.size = Pt(18)
            paragraph.font.color.rgb = colors["text"]
            paragraph.space_before = Pt(4)
    
    def _create_content_slide(self, ppt: Presentation, data: Dict[str, Any], colors: Dict[str, Any]):
        """Create generic content slide"""
        self._create_overview_slide(ppt, data, colors)

class PDFtoPPTSystem:
    """Main system that orchestrates PDF content analysis and PPT creation"""
    
    def __init__(self, aws_region: str = 'us-east-1'):
        self.pdf_extractor = PDFContentExtractor()
        self.content_analyzer = BedrockContentAnalyzer(aws_region)
        self.ppt_generator = PowerPointGenerator()
    
    def generate_ppt_from_pdf(self, pdf_path: str, company_name: str = None, 
                             template: str = "executive") -> str:
        """
        Main function: Generate PowerPoint from PDF content
        
        Args:
            pdf_path: Path to PDF file
            company_name: Company name for customization
            template: "executive", "technical", or "marketing"
            
        Returns:
            Path to generated PowerPoint file
        """
        
        print(f"Generating PowerPoint from PDF content...")
        print(f"PDF: {pdf_path}")
        print(f"Company: {company_name or 'Auto-detect'}")
        print(f"Template: {template}")
        
        # Step 1: Extract content from PDF
        print(f"\n1. Extracting content from PDF...")
        content = self.pdf_extractor.extract_text(pdf_path)
        print(f"   Extracted {len(content):,} characters")
        
        # Step 2: Analyze content with Bedrock
        print(f"\n2. Analyzing content with AWS Bedrock...")
        structure = self.content_analyzer.analyze_content_for_ppt(content, company_name)
        slide_count = len(structure.get('slides', []))
        print(f"   Generated structure for {slide_count} slides")
        
        # Step 3: Create PowerPoint presentation
        print(f"\n3. Creating PowerPoint presentation...")
        ppt_filename = self.ppt_generator.create_presentation(structure, template)
        
        file_size = os.path.getsize(ppt_filename)
        print(f"   Created: {ppt_filename} ({file_size:,} bytes)")
        
        return ppt_filename

def main():
    """User interface for the system"""
    
    print("PowerPoint Generator from PDF Content")
    print("=" * 40)
    
    # Check dependencies
    if not PDF_AVAILABLE:
        print("Missing PDF libraries. Install with:")
        print("pip install PyMuPDF PyPDF2")
        return
    
    if not PPT_AVAILABLE:
        print("Missing PowerPoint library. Install with:")
        print("pip install python-pptx")
        return
    
    # Get inputs
    pdf_path = input("Enter PDF file path: ").strip().strip('"\'')
    
    if not os.path.exists(pdf_path):
        print(f"File not found: {pdf_path}")
        return
    
    company_name = input("Company name (optional): ").strip()
    
    print("\nTemplate options:")
    print("1. Executive (Navy blue)")
    print("2. Technical (Gray-blue)")
    print("3. Marketing (Green)")
    
    choice = input("Choose template (1-3): ").strip()
    templates = {"1": "executive", "2": "technical", "3": "marketing"}
    template = templates.get(choice, "executive")
    
    try:
        # Generate presentation
        system = PDFtoPPTSystem()
        result = system.generate_ppt_from_pdf(pdf_path, company_name, template)
        
        print(f"\nSUCCESS!")
        print(f"Generated: {result}")
        print(f"Location: {os.path.abspath(result)}")
        print(f"Ready to edit in PowerPoint")
        
    except Exception as e:
        print(f"\nError: {e}")
        import traceback
        traceback.print_exc()

if __name__ == "__main__":
    main()